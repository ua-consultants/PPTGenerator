import os
import math
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches

SLIDE_WIDTH = 1920
SLIDE_HEIGHT = 1080
MARGIN = 20

def create_canvas(images):
    n = len(images)
    cols = math.ceil(math.sqrt(n))
    rows = math.ceil(n / cols)

    cell_w = (SLIDE_WIDTH - (cols + 1) * MARGIN) // cols
    cell_h = (SLIDE_HEIGHT - (rows + 1) * MARGIN) // rows

    canvas = Image.new("RGB", (SLIDE_WIDTH, SLIDE_HEIGHT), "white")

    x = y = MARGIN
    index = 0

    for r in range(rows):
        x = MARGIN
        for c in range(cols):
            if index >= n:
                break

            img = Image.open(images[index]).convert("RGB")
            img.thumbnail((cell_w, cell_h))

            paste_x = x + (cell_w - img.width) // 2
            paste_y = y + (cell_h - img.height) // 2

            canvas.paste(img, (paste_x, paste_y))
            x += cell_w + MARGIN
            index += 1

        y += cell_h + MARGIN

    return canvas


def generate_ppt(root_folder, output_path, progress_callback=None):
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    folders = [
        os.path.join(root_folder, d)
        for d in os.listdir(root_folder)
        if os.path.isdir(os.path.join(root_folder, d))
    ]

    total = len(folders)

    for idx, folder in enumerate(folders):
        images = [
            os.path.join(folder, f)
            for f in os.listdir(folder)
            if f.lower().endswith((".jpg", ".jpeg", ".png"))
        ]

        if not images:
            continue

        canvas = create_canvas(images)

        temp_img = "temp_slide.jpg"
        canvas.save(temp_img, quality=85)

        slide = prs.slides.add_slide(blank_layout)

        slide.shapes.add_picture(
            temp_img,
            Inches(0),
            Inches(0),
            width=Inches(13.33)
        )

        title_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(0.1), Inches(12), Inches(0.5)
        )
        title_box.text_frame.text = os.path.basename(folder)

        if progress_callback:
            progress_callback(int((idx + 1) / total * 100))

        os.remove(temp_img)

    prs.save(output_path)
