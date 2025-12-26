# app.py — Slide-by-Slide PPT Generator (with orientation-aware image handling)
import streamlit as st
from pptx import Presentation
from pptx.util import Inches
import os
import tempfile
from io import BytesIO
from PIL import Image as PILImage  # For image dimension handling

# Initialize session state
if 'slides' not in st.session_state:
    st.session_state.slides = []  # List of lists: [[img1, img2], [img3, img4], ...]

st.set_page_config(page_title="🪑 Inspection PPT Builder", layout="wide")
st.title("🪑 Furniture Inspection PPT Builder")
st.markdown("Add images **slide by slide**. When done, generate your report.")

# Show current slides
if st.session_state.slides:
    st.subheader(f"📋 Preview ({len(st.session_state.slides)} slide(s))")
    for i, slide_images in enumerate(st.session_state.slides):
        with st.expander(f"Slide {i+1} ({len(slide_images)} image(s))"):
            cols = st.columns(min(len(slide_images), 4))
            for j, img_path in enumerate(slide_images):
                # Show image preview using path
                cols[j % 4].image(img_path, width=120)

# Main action area
if not st.session_state.slides:
    st.info("👆 Click below to add your first slide.")
    action = "add_first"
else:
    action = st.radio(
        "What would you like to do next?",
        ["Add another slide", "Generate and download PPT"],
        horizontal=True
    )

# Handle actions
if not st.session_state.slides or action == "Add another slide":
    st.subheader("🖼️ Upload images for this slide")
    uploaded_files = st.file_uploader(
        "Drag & drop or browse images (JPG/PNG)",
        type=["jpg", "jpeg", "png"],
        accept_multiple_files=True,
        key="slide_uploader"
    )
    
    if uploaded_files:
        st.success(f"✅ {len(uploaded_files)} image(s) selected for this slide.")
        # Show thumbnails
        cols = st.columns(min(len(uploaded_files), 4))
        for i, f in enumerate(uploaded_files):
            cols[i % 4].image(f, width=100)
        
        col1, col2 = st.columns(2)
        with col1:
            if st.button("✅ Add This Slide"):
                # Save to temp dir and store paths
                temp_dir = tempfile.mkdtemp()
                image_paths = []
                for f in uploaded_files:
                    path = os.path.join(temp_dir, f.name)
                    with open(path, "wb") as fp:
                        fp.write(f.getvalue())
                    image_paths.append(path)
                st.session_state.slides.append(image_paths)
                st.rerun()
        with col2:
            st.button("↩️ Cancel")

# Generate PPT
if st.session_state.slides and action == "Generate and download PPT":
    if st.button("📥 Generate and Download PPT", type="primary"):
        if not st.session_state.slides:
            st.error("❌ No slides to generate!")
        else:
            # Create PPT
            prs = Presentation()
            # Find blank layout
            blank_layout = None
            for layout in prs.slide_layouts:
                if "blank" in layout.name.lower():
                    blank_layout = layout
                    break
            if blank_layout is None:
                blank_layout = prs.slide_layouts[0]

            for slide_images in st.session_state.slides:
                slide = prs.slides.add_slide(blank_layout)
                num_imgs = len(slide_images)
                
                # Auto grid based on number of images
                if num_imgs <= 4:
                    rows, cols = (2, 2)
                elif num_imgs <= 6:
                    rows, cols = (2, 3)
                elif num_imgs <= 9:
                    rows, cols = (3, 3)
                else:
                    rows, cols = (4, 4)
                
                # Slide usable area (in inches)
                max_width = 15.0   # Approx full slide width
                max_height = 7.5   # Leave room for title/footer if needed
                img_width = max_width / cols
                img_height = max_height / rows

                for idx, img_path in enumerate(slide_images):
                    row = idx // cols
                    col = idx % cols

                    # Cell boundaries (in inches)
                    cell_left_in = col * img_width + 0.15
                    cell_top_in = row * img_height + 0.3
                    cell_width_in = img_width - 0.3
                    cell_height_in = img_height - 0.3

                    # Load image to get native dimensions (do NOT auto-rotate)
                    try:
                        with PILImage.open(img_path) as pil_img:
                            # Make a copy to avoid "closed file" issues
                            pil_img.load()  # Ensure pixels are loaded
                            orig_width_px, orig_height_px = pil_img.size

                            # Get DPI or default to 96
                            dpi = pil_img.info.get('dpi', (96, 96))
                            dpi_x = dpi[0] if isinstance(dpi, tuple) else dpi
                            if dpi_x <= 1:
                                dpi_x = 96

                            # Convert pixels to inches
                            orig_width_in = orig_width_px / dpi_x
                            orig_height_in = orig_height_px / dpi_x

                        # Compute scale to fit image inside cell while preserving aspect ratio
                        scale_x = cell_width_in / orig_width_in
                        scale_y = cell_height_in / orig_height_in
                        scale = min(scale_x, scale_y)

                        new_width_in = orig_width_in * scale
                        new_height_in = orig_height_in * scale

                        # Center in cell
                        centered_left_in = cell_left_in + (cell_width_in - new_width_in) / 2
                        centered_top_in = cell_top_in + (cell_height_in - new_height_in) / 2

                        # Add to slide
                        slide.shapes.add_picture(
                            img_path,
                            Inches(centered_left_in),
                            Inches(centered_top_in),
                            width=Inches(new_width_in),
                            height=Inches(new_height_in)
                        )
                    except Exception as e:
                        st.warning(f"⚠️ Failed to add image: {os.path.basename(img_path)} – {str(e)}")

            # Save and download
            ppt_buffer = BytesIO()
            prs.save(ppt_buffer)
            ppt_buffer.seek(0)

            st.download_button(
                label="⬇️ Download Inspection Report (PPTX)",
                data=ppt_buffer,
                file_name="Furniture_Inspection_Report.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

# Reset button
if st.session_state.slides:
    if st.button("🔄 Start Over"):
        st.session_state.slides = []
        st.rerun()
